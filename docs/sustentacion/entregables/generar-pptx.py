"""
Genera SaludDeUna-Sustentacion-Tecnica.pptx
Requiere: pip install python-pptx
Ejecutar desde la carpeta entregables/:  python generar-pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pathlib import Path
import os

# ── Paths ──────────────────────────────────────────────────────────────────
BASE = Path(__file__).parent
ASSETS = BASE / "presentacion-web" / "assets"
OUT = BASE / "SaludDeUna-Sustentacion-Tecnica.pptx"

# ── Design System ───────────────────────────────────────────────────────────
BG        = RGBColor(0x02, 0x08, 0x17)   # #020817
PANEL     = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A
PANEL2    = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B
BLUE      = RGBColor(0x3B, 0x82, 0xF6)   # #3B82F6
MINT      = RGBColor(0x10, 0xB9, 0x81)   # #10B981
VIOLET    = RGBColor(0x8B, 0x5C, 0xF6)   # #8B5CF6
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # #F59E0B
CORAL     = RGBColor(0xFB, 0x71, 0x85)   # #FB7185
TEXT      = RGBColor(0xF1, 0xF5, 0xF9)   # #F1F5F9
MUTED     = RGBColor(0x94, 0xA3, 0xB8)   # #94A3B8
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# Slide size: 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank

# ── Helpers ─────────────────────────────────────────────────────────────────

def new_slide():
    slide = prs.slides.add_slide(blank)
    # dark background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def box(slide, x, y, w, h, bg=None, border=None):
    from pptx.util import Pt
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    if bg:
        shp.fill.solid()
        shp.fill.fore_color.rgb = bg
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, x, y, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
        italic=False, wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or TEXT
    run.font.name = font_name
    return txb

def heading(slide, text, size=36, color=None):
    txt(slide, text, 0.5, 0.75, 12.3, 1.5, size=size, bold=True, color=color or TEXT)

def subheading(slide, text):
    txt(slide, text, 0.5, 1.7, 12.3, 0.6, size=14, color=MUTED)

def tag(slide, section_text, slide_num, total=25):
    txt(slide, section_text.upper(), 0.5, 0.15, 8, 0.35, size=10, bold=True,
        color=MINT, font_name="Courier New")
    txt(slide, f"{slide_num:02d} / {total:02d}", 11.5, 0.15, 1.3, 0.35, size=10, bold=True,
        color=MUTED, align=PP_ALIGN.RIGHT, font_name="Courier New")

def qr_badge(slide):
    qr_path = ASSETS / "qr-apk.png"
    if qr_path.exists():
        slide.shapes.add_picture(str(qr_path), Inches(12.5), Inches(6.8), Inches(0.7), Inches(0.7))

def card(slide, x, y, w, h, title, body, title_color=None):
    box(slide, x, y, w, h, bg=PANEL2, border=RGBColor(0x2D,0x3F,0x56))
    if title:
        txt(slide, title, x+0.15, y+0.12, w-0.3, 0.35, size=11, bold=True,
            color=title_color or MINT)
    if body:
        txt(slide, body, x+0.15, y+0.5, w-0.3, h-0.65, size=10, color=TEXT)

def metric_box(slide, x, y, w, h, value, label, val_color=None):
    box(slide, x, y, w, h, bg=PANEL, border=RGBColor(0x2D,0x3F,0x56))
    txt(slide, value, x, y+0.1, w, 0.6, size=32, bold=True,
        color=val_color or BLUE, align=PP_ALIGN.CENTER)
    txt(slide, label.upper(), x, y+0.7, w, 0.35, size=8, bold=True,
        color=MUTED, align=PP_ALIGN.CENTER, font_name="Courier New")

def img(slide, filename, x, y, w, h):
    p = ASSETS / filename
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        # placeholder
        b = box(slide, x, y, w, h, bg=PANEL2, border=MUTED)

def bullet_list(slide, items, x, y, w, h, size=12, color=None):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"→  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color or TEXT
        run.font.name = "Calibri"

def flow_strip(slide, steps, y=3.8, colors=None):
    n = len(steps)
    step_w = 10.5 / n
    x_start = 0.5
    for i, step in enumerate(steps):
        x = x_start + i * step_w
        c = colors[i] if colors else PANEL2
        box(slide, x, y, step_w - 0.1, 0.65, bg=c, border=BLUE)
        txt(slide, step, x+0.05, y+0.05, step_w-0.2, 0.55,
            size=9, bold=True, align=PP_ALIGN.CENTER, color=TEXT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 · PORTADA
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
# gradient boxes for visual depth
box(s, 0, 0, 13.33, 7.5, bg=RGBColor(0x07,0x11,0x1F))
box(s, 8, 0, 5.33, 7.5, bg=RGBColor(0x0D,0x1B,0x2F))

logo_path = ASSETS / "logoSaludDeUna.png"
if logo_path.exists():
    s.shapes.add_picture(str(logo_path), Inches(0.5), Inches(0.5), Inches(2.5), Inches(0.9))

txt(s, "SUSTENTACIÓN TÉCNICA · IETI 2026-1 · 21 MAYO 2026",
    0.5, 1.6, 7, 0.4, size=10, bold=True, color=MINT, font_name="Courier New")
txt(s, "SaludDeUna", 0.5, 2.0, 9, 1.8, size=72, bold=True, color=TEXT)
txt(s, "Plataforma de Comunicación Clínica Asistida por IA",
    0.5, 3.8, 8.5, 0.6, size=18, color=MUTED)

# tech chips
chips = ["NestJS 11", "Expo 55", "Next.js 16", "Gemini 2.5", "MongoDB", "Redis", "Auth0"]
for i, chip in enumerate(chips):
    cx = 0.5 + i * 1.8
    if cx > 12: break
    box(s, cx, 4.6, 1.6, 0.38, bg=PANEL2, border=BLUE)
    txt(s, chip, cx+0.05, 4.63, 1.5, 0.32, size=9, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

# QR panel
box(s, 10.2, 1.2, 2.7, 3.8, bg=WHITE, border=None)
qr_path = ASSETS / "qr-apk.png"
if qr_path.exists():
    s.shapes.add_picture(str(qr_path), Inches(10.5), Inches(1.5), Inches(2.1), Inches(2.1))
txt(s, "APK disponible", 10.2, 3.7, 2.7, 0.4, size=11, bold=True,
    color=RGBColor(0x0F,0x17,0x2A), align=PP_ALIGN.CENTER)
txt(s, "Escanear para instalar demo", 10.2, 4.1, 2.7, 0.35, size=9,
    color=RGBColor(0x47,0x55,0x69), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 · AGENDA
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "AGENDA", 2)
heading(s, "El recorrido demuestra problema, solución y solidez técnica.")
subheading(s, "7 bloques temáticos · Demo en vivo como evidencia final · ~20 minutos")

agenda = [
    ("01", "Problema\ne Impacto", "3 min"),
    ("02", "Propuesta\nde Valor", "2 min"),
    ("03", "MVP y\nCasos de Uso", "2.5 min"),
    ("04", "Arquitectura\nC4", "3 min"),
    ("05", "Calidad del\nSoftware", "3 min"),
    ("06", "IA en Producto\ny Desarrollo", "2.5 min"),
    ("07", "Demo +\nConclusiones", "5–7 min"),
]
w = 1.7
for i, (num, name, time) in enumerate(agenda):
    x = 0.3 + i * 1.83
    box(s, x, 2.5, w, 2.8, bg=PANEL, border=RGBColor(0x2D,0x3F,0x56))
    txt(s, num, x, 2.6, w, 0.6, size=28, bold=True, color=BLUE,
        align=PP_ALIGN.CENTER, font_name="Courier New")
    txt(s, name, x, 3.2, w, 0.9, size=10, bold=True, color=TEXT,
        align=PP_ALIGN.CENTER)
    txt(s, time, x, 4.1, w, 0.35, size=9, color=MUTED,
        align=PP_ALIGN.CENTER, font_name="Courier New")
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 · PROBLEMA EN UNA FRASE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "CONTEXTO DEL PROBLEMA", 3)
txt(s, '"La información clínica llega al médico\nmal organizada, tarde, o incompleta."',
    0.7, 1.0, 12, 2.5, size=34, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
txt(s, "Pacientes de atención primaria en Colombia — Consulta ambulatoria",
    0.7, 3.6, 12, 0.45, size=13, color=MUTED, align=PP_ALIGN.CENTER)

stats = [("10 min", "Duración consulta promedio"), ("40–60%", "Tiempo en anamnesis"), ("Alta carga", "Triage manual en urgencias")]
for i, (val, lbl) in enumerate(stats):
    x = 1.5 + i * 3.5
    box(s, x, 4.3, 3.0, 1.5, bg=PANEL2, border=BLUE)
    txt(s, val, x, 4.4, 3.0, 0.65, size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, lbl, x, 5.1, 3.0, 0.5, size=10, color=MUTED, align=PP_ALIGN.CENTER)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 · LADO DEL PACIENTE
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "PROBLEMA · PACIENTE", 4)
heading(s, "El paciente no sabe comunicar síntomas de forma clínicamente útil.")
subheading(s, "El problema no es la información — es la falta de estructura para transmitirla.")

card(s, 0.4, 2.4, 6.0, 3.4, "❌  Dolor de comunicación", None, CORAL)
txt(s, "No distingue si el síntoma es urgente o puede esperar",
    0.6, 2.9, 5.7, 0.45, size=12, color=TEXT)
txt(s, "Olvida antecedentes, medicamentos actuales y duración exacta",
    0.6, 3.4, 5.7, 0.45, size=12, color=TEXT)
txt(s, "Describe síntomas sin contexto temporal ni evolutivo claro",
    0.6, 3.9, 5.7, 0.45, size=12, color=TEXT)
txt(s, "Consultas cortas y costosas con información fragmentada",
    0.6, 4.4, 5.7, 0.45, size=12, color=TEXT)

card(s, 6.8, 2.4, 6.0, 3.4, "⚠️  Consecuencias reales", None, AMBER)
txt(s, "Visitas innecesarias a urgencias por falta de triage previo",
    7.0, 2.9, 5.7, 0.45, size=12, color=TEXT)
txt(s, "Diagnósticos tardíos por información incompleta al médico",
    7.0, 3.4, 5.7, 0.45, size=12, color=TEXT)
txt(s, "Sistema saturado de urgencias por casos no urgentes",
    7.0, 3.9, 5.7, 0.45, size=12, color=TEXT)
txt(s, "Sin continuidad estructurada entre consultas del mismo paciente",
    7.0, 4.4, 5.7, 0.45, size=12, color=TEXT)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 · LADO DEL MÉDICO
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "PROBLEMA · MÉDICO", 5)
heading(s, "El médico reconstruye el contexto clínico desde cero en cada consulta.")
subheading(s, "40–60% del tiempo clínico en anamnesis repetida — información que el paciente ya tiene.")

problems = [
    ("Anamnesis repetida", "El médico hace las mismas preguntas básicas en cada consulta. Información que el paciente ya tiene pero no estructura.", AMBER),
    ("Sin continuidad", "Entre una consulta y la siguiente no hay contexto estructurado. El historial existe pero fragmentado.", AMBER),
    ("Alta carga operativa", "El triage manual consume tiempo de atención clínica. El médico es el cuello de botella del sistema.", CORAL),
]
for i, (title, body, color) in enumerate(problems):
    x = 0.4 + i * 4.3
    card(s, x, 2.5, 4.1, 2.8, title, body, color)

box(s, 0.4, 5.5, 12.5, 0.6, bg=RGBColor(0x10,0x3D,0x2B), border=MINT)
txt(s, "40–60% de cada consulta se va en recolectar información que el paciente ya tiene.",
    0.6, 5.55, 12.2, 0.5, size=12, bold=True, color=RGBColor(0xD1,0xFA,0xE5), align=PP_ALIGN.CENTER)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 · PROPUESTA DE VALOR
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "PROPUESTA DE VALOR", 6)
heading(s, "SaludDeUna es la capa inteligente entre paciente y médico.")
subheading(s, "IA con guardrails clínicos que estructura información, clasifica prioridad y da continuidad.")

pillars = [
    ("01 · Triage Asistido IA", "Cuestionario estructurado por especialidad. Detección de red flags. Clasificación de prioridad LOW / MODERATE / HIGH antes de que el médico vea al paciente.", BLUE),
    ("02 · Resumen Clínico", "Gemini 2.5-flash genera nota médica estructurada basada en el cuestionario y el chat. El médico llega informado, no en blanco.", MINT),
    ("03 · Seguimiento Post-consulta", "Followups automáticos a 72 horas y 7 días. Si el paciente empeora ≥2 puntos, el sistema escala creando una nueva consulta de alta prioridad.", VIOLET),
]
for i, (title, body, color) in enumerate(pillars):
    x = 0.4 + i * 4.3
    box(s, x, 2.4, 4.1, 3.0, bg=PANEL2, border=color)
    txt(s, title, x+0.15, 2.5, 3.8, 0.4, size=11, bold=True, color=color)
    txt(s, body, x+0.15, 2.95, 3.8, 2.35, size=11, color=TEXT)

box(s, 0.4, 5.65, 12.5, 0.55, bg=RGBColor(0x10,0x3D,0x2B), border=MINT)
txt(s, "La IA organiza, clasifica y comunica. El médico decide.",
    0.6, 5.7, 12.2, 0.45, size=13, bold=True, color=RGBColor(0xD1,0xFA,0xE5), align=PP_ALIGN.CENTER)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 · QUÉ NO ES
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "ALCANCE Y ÉTICA", 7)
heading(s, "Claridad de alcance: SaludDeUna NO reemplaza al médico.")
subheading(s, "Límites técnicos y éticos explícitos — no una promesa, sino código ejecutable.")

box(s, 0.4, 2.4, 5.9, 4.1, bg=RGBColor(0x2D,0x08,0x0E), border=CORAL)
txt(s, "❌  Fuera de alcance", 0.6, 2.5, 5.6, 0.4, size=12, bold=True, color=CORAL)
noscope = ["Telemedicina regulada (Res. 2654/2019)", "Diagnóstico automático de enfermedades",
           "Prescripción de medicamentos", "Integración productiva con HCE de terceros",
           "Videollamada en tiempo real", "Cobro productivo real"]
for i, item in enumerate(noscope):
    txt(s, f"×  {item}", 0.6, 3.0 + i*0.5, 5.6, 0.45, size=11, color=TEXT)

box(s, 6.8, 2.4, 6.1, 4.1, bg=RGBColor(0x06,0x28,0x20), border=MINT)
txt(s, "✅  Lo que SÍ hace la IA", 7.0, 2.5, 5.8, 0.4, size=12, bold=True, color=MINT)
inscope = ["Clasifica prioridad clínica (LOW/MODERATE/HIGH)", "Detecta red flags en las respuestas del paciente",
           "Genera resumen clínico neutral de urgencia", "Enriquece contexto con evidencia clínica (RAG)",
           "Guardrail bloquea diagnóstico y prescripción", "Fallback a motor de reglas si Gemini falla"]
for i, item in enumerate(inscope):
    txt(s, f"→  {item}", 7.0, 3.0 + i*0.5, 5.8, 0.45, size=11, color=TEXT)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 · ALCANCE MVP
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "MVP IMPLEMENTADO", 8)
heading(s, "Todo lo que se prometió está implementado y con tests.")
subheading(s, "3 especialidades · 3 roles · 2 frontends · 17 módulos backend · 406+ tests pasando")

mvp = [
    ("Auth & Roles", "Registro, login JWT+Auth0, RBAC, máx 3 sesiones, refresh tokens"),
    ("Triage IA", "Cuestionario por especialidad, red flags, prioridad, guardrail, fallback"),
    ("Cola Médica", "PENDING→IN_PROGRESS→COMPLETED, ordenada por prioridad"),
    ("Chat Real-time", "Socket.IO, persistido en MongoDB, Redis adapter multi-instancia"),
    ("Resumen IA", "Gemini genera nota clínica, guardrail valida, historial acumulativo"),
    ("Seguimiento", "Followups 72h + 7d, escalación automática si empeora, BullMQ"),
    ("Admin Panel", "Verificación REThUS, gestión usuarios, billing, dashboard KPIs"),
    ("Billing", "Precios por especialidad COP, transacciones, revenue metrics"),
]
cols = 4
for i, (title, body) in enumerate(mvp):
    row, col = divmod(i, cols)
    x = 0.35 + col * 3.25
    y = 2.5 + row * 1.85
    card(s, x, y, 3.1, 1.7, title, body)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 · TRES ROLES
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "CASOS DE USO · ROLES", 9)
heading(s, "Tres roles, interfaces separadas, backend compartido.")
subheading(s, "Patient-only en mobile · Doctor y Admin en web · RBAC en cada endpoint")

roles_data = [
    ("Paciente", "EXPO 55 · REACT NATIVE", BLUE,
     ["Registro y autenticación", "Selección de especialidad", "Triage IA guiado",
      "Checkout y pago simulado", "Chat con médico en tiempo real", "Seguimiento post-consulta", "Historial clínico acumulativo"]),
    ("Médico", "NEXT.JS 16 · APP ROUTER", MINT,
     ["Registro y verificación REThUS", "Cola de consultas por prioridad",
      "Chat bidireccional con paciente", "Generación resumen clínico IA",
      "Cierre y calificación consulta", "Dashboard de métricas propias"]),
    ("Administrador", "NEXT.JS 16 · APP ROUTER", VIOLET,
     ["Verificación manual REThUS", "Gestión de usuarios y roles",
      "Precios por especialidad", "Dashboard KPIs del negocio",
      "Revenue metrics y transacciones", "Métricas técnicas plataforma"]),
]
for i, (role, platform, color, features) in enumerate(roles_data):
    x = 0.35 + i * 4.3
    box(s, x, 2.4, 4.1, 4.6, bg=PANEL2, border=color)
    txt(s, role, x+0.15, 2.5, 3.8, 0.45, size=18, bold=True, color=color)
    txt(s, platform, x+0.15, 2.95, 3.8, 0.3, size=9, bold=True, color=MUTED, font_name="Courier New")
    for j, feat in enumerate(features):
        txt(s, f"  {feat}", x+0.15, 3.35 + j*0.47, 3.8, 0.42, size=10, color=TEXT)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 · FLUJO E2E
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "FLUJO DEL NEGOCIO E2E", 10)
heading(s, "Flujo completo de principio a fin — trazable y auditable.")
subheading(s, "Cada consulta tiene CorrelationId. Cada evento pasa por el Outbox transaccional.")

steps = ["🔐 Login\nJWT+Auth0", "🤖 Triage IA\nGemini+Guard", "💳 Checkout\nBilling COP",
         "📋 Cola\nPor prioridad", "💬 Chat\nSocket.IO", "📝 Resumen\nGemini", "🔔 Followup\n72h + 7d"]
for i, step in enumerate(steps):
    x = 0.4 + i * 1.8
    box(s, x, 2.6, 1.65, 1.0, bg=PANEL2, border=BLUE)
    txt(s, step, x+0.05, 2.65, 1.55, 0.9, size=9, bold=True, align=PP_ALIGN.CENTER, color=TEXT)
    if i < len(steps)-1:
        txt(s, "→", x+1.55, 2.9, 0.3, 0.4, size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

details = [
    ("Mensajes persistidos en MongoDB con correlationId", BLUE),
    ("Outbox transaccional crea eventos de dominio", MINT),
    ("Worker BullMQ procesa followups asincrónicamente", VIOLET),
    ("Escalación automática si severidad aumenta ≥2 puntos", AMBER),
]
for i, (detail, color) in enumerate(details):
    x = 0.4 + i * 3.25
    box(s, x, 4.0, 3.1, 0.85, bg=PANEL, border=color)
    txt(s, detail, x+0.1, 4.07, 2.9, 0.75, size=10, color=TEXT)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 · C4 NIVEL 1
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "ARQUITECTURA · C4 NIVEL 1", 11)
heading(s, "Diagrama de Contexto: 3 actores · 1 sistema central · 6 servicios externos.")
subheading(s, "Todo HTTPS · Chat por WebSocket vía Socket.IO")
img(s, "sad_c4_context_diagram.png", 0.5, 2.2, 12.33, 4.7)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 · C4 NIVEL 2
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "ARQUITECTURA · C4 NIVEL 2", 12)
heading(s, "Contenedores: frontends, API, Worker y capa de datos.")
subheading(s, "API para request-response · Worker para jobs asincrónicos · MongoDB Vector Search 768d")
img(s, "sad_c4_container_view.png", 0.5, 2.2, 12.33, 4.7)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 · BACKEND 17 MÓDULOS
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "BACKEND · 17 MÓDULOS NESTJS", 13)
heading(s, "4 dominios · 17 módulos · 3 capas transversales.")
subheading(s, "Cada módulo tiene límites claros. Acoplamiento por DI, nunca por imports directos.")

domains = [
    ("🟢 USUARIOS", ["Auth (JWT+Auth0)", "Patients", "Doctors (REThUS)", "Admins"], MINT),
    ("🟡 CLÍNICO", ["Triage + AI", "Consultations", "Chat (Socket.IO)", "Billing · Followups"], AMBER),
    ("🟣 IA / RAG", ["AiModule (Gemini)", "Knowledge (Chunks)", "Rag (8 pasos)", "Vector Search 768d"], VIOLET),
    ("🔴 INFRA", ["Outbox (exactly-once)", "Dashboard (KPIs)", "Common (Guards)", "OTel Bootstrap"], CORAL),
]
for i, (domain, items, color) in enumerate(domains):
    x = 0.4 + i * 3.2
    box(s, x, 2.4, 3.05, 3.5, bg=PANEL, border=color)
    txt(s, domain, x+0.1, 2.5, 2.85, 0.4, size=10, bold=True, color=color, font_name="Courier New")
    for j, item in enumerate(items):
        box(s, x+0.1, 3.0 + j*0.7, 2.85, 0.6, bg=RGBColor(0x07,0x11,0x1F), border=RGBColor(0x2D,0x3F,0x56))
        txt(s, item, x+0.2, 3.08 + j*0.7, 2.65, 0.5, size=10, color=TEXT)

box(s, 0.4, 6.1, 12.5, 0.55, bg=RGBColor(0x07,0x1A,0x2E), border=BLUE)
txt(s, "JWT + RBAC  ·  Throttling 20 req/60s  ·  GlobalExceptionFilter + CorrelationId",
    0.6, 6.15, 12.2, 0.45, size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER, font_name="Courier New")
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 · PIPELINE TRIAGE IA
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "IA · PIPELINE DE TRIAGE", 14)
heading(s, "Pipeline: cuestionario → Gemini → guardrail → resultado auditable.")
subheading(s, "Fallback a RedFlagsEngine si Gemini falla — sin error visible para el usuario.")

pipe = [
    ("📋 Cuestionario", "5–8 preguntas\npor especialidad", BLUE),
    ("🔍 RAG Context", "Vector Search 768d\nRedis cache 300s", MINT),
    ("🤖 Gemini 2.5-flash", "Prompt versionado\nAudit log", VIOLET),
    ("🛡️ Guardrail", "Bloquea diagnóstico\ny prescripción", CORAL),
    ("✅ Resultado", "LOW/MODERATE/HIGH\n+ red flags", MINT),
]
for i, (name, desc, color) in enumerate(pipe):
    x = 0.4 + i * 2.5
    box(s, x, 2.5, 2.35, 2.0, bg=PANEL2, border=color)
    txt(s, name, x+0.1, 2.58, 2.15, 0.5, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, desc, x+0.1, 3.1, 2.15, 1.2, size=10, color=TEXT, align=PP_ALIGN.CENTER)
    if i < len(pipe)-1:
        txt(s, "→", x+2.25, 3.2, 0.35, 0.5, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

box(s, 0.4, 4.8, 5.9, 1.3, bg=RGBColor(0x2D,0x08,0x0E), border=CORAL)
txt(s, "Guardrail activa: aiSummary = null · guardrailApplied = true · log WARN + correlationId",
    0.6, 4.88, 5.7, 0.55, size=10, bold=True, color=CORAL)
txt(s, "La respuesta de IA no se persiste si hay diagnóstico o prescripción.",
    0.6, 5.45, 5.7, 0.55, size=10, color=TEXT)

box(s, 6.8, 4.8, 5.9, 1.3, bg=RGBColor(0x2D,0x1A,0x00), border=AMBER)
txt(s, "Fallback: Gemini falla → RedFlagsEngine (reglas locales) → analysisType = RULE_BASED",
    7.0, 4.88, 5.7, 0.55, size=10, bold=True, color=AMBER)
txt(s, "Sin error visible para el usuario. Confiabilidad garantizada.",
    7.0, 5.45, 5.7, 0.55, size=10, color=TEXT)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 · CRITERIOS DE CALIDAD
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "CRITERIOS DE CALIDAD", 15)
heading(s, "Seis criterios de calidad con evidencia técnica medible.")
subheading(s, "No son aspiraciones — son mecanismos implementados con tests que los validan.")

criteria = [
    ("01", "🔐 Seguridad", "JWT HS256+RS256, RBAC, rate limiting Redis, bcrypt cost 12, máx 3 sesiones", BLUE),
    ("02", "✅ Confiabilidad", "406+ tests · 93% coverage · E2E con MongoDB real · CI falla si <80%", MINT),
    ("03", "📊 Observabilidad", "OpenTelemetry · logs JSON estructurados · correlationId · p95 latencia · KPIs", VIOLET),
    ("04", "⚡ Rendimiento", "Redis cache 300s · Rate limiting 20 req/60s · SLO p95 <1500ms · BullMQ async", AMBER),
    ("05", "🏗️ Mantenibilidad", "17 módulos NestJS · Clean Architecture · DI explícita · Swagger autogenerado", CORAL),
    ("06", "📈 Escalabilidad", "Redis adapter Socket.IO · BullMQ multi-worker · ECS Fargate autoscaling", BLUE),
]
for i, (num, name, desc, color) in enumerate(criteria):
    row, col = divmod(i, 3)
    x = 0.35 + col * 4.3
    y = 2.5 + row * 2.1
    box(s, x, y, 4.1, 1.9, bg=PANEL, border=color)
    txt(s, num, x+0.15, y+0.1, 0.6, 0.5, size=22, bold=True, color=color, font_name="Courier New")
    txt(s, name, x+0.8, y+0.1, 3.15, 0.4, size=12, bold=True, color=TEXT)
    txt(s, desc, x+0.15, y+0.6, 3.8, 1.1, size=10, color=MUTED)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 · SEGURIDAD
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "CALIDAD · SEGURIDAD", 16)
heading(s, "Defensa en profundidad antes del controlador de negocio.")
subheading(s, "Cada request pasa por 5 capas de validación antes de ejecutar lógica de dominio.")

sec = [
    ("🔑  Dual Auth: JWT + Auth0", "HS256 para flujo legacy · RS256 vía JWKS para OAuth2/Auth0\nJwtAuthGuard global con bypass @Public() por excepción", BLUE),
    ("👥  RBAC Explícito", "3 roles: PATIENT · DOCTOR · ADMIN · Guards globales via APP_GUARD\nDoctorVerifiedGuard en todas las rutas clínicas", MINT),
    ("🚦  Rate Limiting Distribuido", "20 requests / 60 segundos / IP · Redis Store distribuido\nFunciona en multi-instancia sin acoplar estado local", VIOLET),
    ("🔒  Gestión de Sesiones", "bcrypt cost 12 · Máximo 3 sesiones activas por usuario\nLIFO revocation automática · Refresh token rotation", AMBER),
]
for i, (title, body, color) in enumerate(sec):
    row, col = divmod(i, 2)
    x = 0.35 + col * 6.5
    y = 2.5 + row * 2.0
    box(s, x, y, 6.2, 1.8, bg=PANEL2, border=color)
    txt(s, title, x+0.15, y+0.1, 5.9, 0.4, size=12, bold=True, color=color)
    txt(s, body, x+0.15, y+0.6, 5.9, 1.05, size=11, color=TEXT)

box(s, 0.35, 6.55, 12.6, 0.55, bg=RGBColor(0x07,0x1A,0x2E), border=BLUE)
txt(s, "Cada error de API incluye statusCode, message y correlationId — auditable y trazable.",
    0.55, 6.6, 12.3, 0.45, size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 · TESTING
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "CALIDAD · CONFIABILIDAD", 17)
heading(s, "406 tests, 93% cobertura, E2E con MongoDB real en memoria.")
subheading(s, "Si la cobertura baja del 80%, el CI falla automáticamente. No hay degradación silenciosa.")

metric_box(s, 0.4, 2.5, 2.9, 1.5, "406+", "Tests backend", BLUE)
metric_box(s, 3.6, 2.5, 2.9, 1.5, "93%", "Stmt coverage", MINT)
metric_box(s, 6.8, 2.5, 2.9, 1.5, "80%", "Umbral mínimo CI", VIOLET)
metric_box(s, 10.0, 2.5, 2.9, 1.5, "82", "Tests web portal", AMBER)

card(s, 0.4, 4.3, 6.0, 2.3,
     "Unit Tests (Jest + NestJS Testing)",
     "Mocks de módulos NestJS · AiService siempre mockeado · Validan lógica de dominio aislada\nRápidos, ejecutan en CI en <30s · Cero dependencias externas", BLUE)
card(s, 6.8, 4.3, 6.0, 2.3,
     "E2E Tests (mongodb-memory-server)",
     "MongoDB real en memoria para cada suite · No se mockea la capa de BD\nDetectan bugs de queries complejas y relaciones entre esquemas · 12 suites E2E", MINT)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 · OBSERVABILIDAD
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "CALIDAD · OBSERVABILIDAD", 18)
heading(s, "Tres niveles: logs estructurados, métricas técnicas y trazas distribuidas.")
subheading(s, "Cada request queda trazado. Cada error tiene correlationId. Cada KPI disponible via API.")

obs = [
    ("📋  Logs Estructurados", 'JSON con: correlationId · userId · method · path\nstatusCode · durationMs — cada request trazable', BLUE),
    ("📊  Métricas Técnicas", 'p95 latencia en tiempo real · Error rate\nEndpoint /v1/dashboard/technical · Alerta si p95 > 1500ms', MINT),
    ("🔭  OpenTelemetry", 'NodeSDK con auto-instrumentación · OTLP HTTP :4318\nListo para Jaeger/Tempo en producción', VIOLET),
]
for i, (title, body, color) in enumerate(obs):
    x = 0.35 + i * 4.3
    card(s, x, 2.5, 4.1, 2.3, title, body, color)

kpis = ["Tiempo a 1ª respuesta médica", "Tasa uso resumen clínico IA", "Tasa confirmación red flags", "Retención followup 7 días"]
for i, kpi in enumerate(kpis):
    x = 0.35 + i * 3.25
    box(s, x, 5.1, 3.1, 0.85, bg=PANEL, border=BLUE)
    txt(s, kpi, x+0.1, 5.2, 2.9, 0.65, size=10, bold=True, color=BLUE, align=PP_ALIGN.CENTER, font_name="Courier New")
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 · IA EN EL PRODUCTO
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "IA · EN EL PRODUCTO", 19)
heading(s, "La IA tiene tres funciones concretas, todas con guardrails.")
subheading(s, "Gemini 2.5-flash + embedding-001 · RAG con Vector Search HNSW 768 dimensiones")

ia = [
    ("🤖  Triage Inteligente", "Clasifica prioridad LOW/MODERATE/HIGH\nDetecta red flags por especialidad · Prompt versionado\nAudit log de cada análisis · Fallback a ReglaEngine", VIOLET),
    ("📝  Resumen Clínico", "Nota médica estructurada post-chat\nDoctor solicita al cerrar consulta\nGuardrail valida antes de mostrar\nCampo aiSummary en Consultation schema", VIOLET),
    ("🔍  RAG Pipeline", "Retrieval-Augmented Generation — 8 pasos\nKnowledge base clínica con chunking\nVector Search Atlas HNSW · Cache Redis 300s\nRagTrace auditable en MongoDB", VIOLET),
]
for i, (title, body, color) in enumerate(ia):
    x = 0.35 + i * 4.3
    card(s, x, 2.5, 4.1, 3.2, title, body, color)

img(s, "sad_rag_pipeline_diagram.png", 0.35, 5.9, 12.6, 1.2)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 · GUARDRAIL CLÍNICO
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "IA · GUARDRAIL CLÍNICO", 20)
heading(s, "El guardrail no es una promesa — es código ejecutable y auditable.")
subheading(s, "GuardrailService evalúa cada salida de Gemini antes de persistir. Sin excepciones.")

box(s, 0.4, 2.4, 12.5, 2.8, bg=RGBColor(0x1E,0x0D,0x3B), border=VIOLET)
txt(s, "Categorías bloqueadas por GuardrailService:", 0.6, 2.5, 12.2, 0.4, size=12, bold=True, color=VIOLET)
rules = [
    "🚫  Lenguaje de diagnóstico:  'Usted tiene', 'padece de', 'el diagnóstico es', 'parece ser X enfermedad'",
    "🚫  Lenguaje de prescripción:  'Tome', 'debe tomar', 'le recomiendo X medicamento', 'dosis de X'",
    "🚫  Afirmaciones clínicas definitivas:  'Este síntoma confirma', 'con certeza indica', 'definitivamente es'",
]
for i, rule in enumerate(rules):
    txt(s, rule, 0.6, 3.0 + i*0.65, 12.2, 0.6, size=11, color=TEXT)

for i, (val, lbl, color) in enumerate([
    ("aiSummary = null", "La respuesta de IA no se persiste", CORAL),
    ("guardrailApplied = true", "Flag auditable en el documento", AMBER),
    ("log WARN + correlationId", "Evento registrado y trazable", BLUE),
]):
    x = 0.4 + i * 4.35
    box(s, x, 5.55, 4.1, 1.3, bg=PANEL2, border=color)
    txt(s, val, x+0.1, 5.62, 3.9, 0.5, size=12, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Courier New")
    txt(s, lbl, x+0.1, 6.15, 3.9, 0.5, size=10, color=TEXT, align=PP_ALIGN.CENTER)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 · IA EN EL DESARROLLO
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "IA · EN EL DESARROLLO (SWNT §6)", 21)
heading(s, "Claude Code fue el copiloto técnico del equipo durante todo el proyecto.")
subheading(s, "Evidencia concreta: generación de código, tests, arquitectura y documentación.")

dev_ia = [
    ("⚙️  Generación de Código", "Módulos NestJS completos · Controllers, Services, DTOs\nSchemas Mongoose · Hooks React y componentes Next.js", VIOLET),
    ("🧪  Testing", "406 tests generados/revisados con IA · E2E suites\nmongodb-memory-server · Mocks de módulos NestJS", VIOLET),
    ("🏛️  Arquitectura", "Revisiones arquitectónicas por sprint · Auditoría C4\nDiagrams as code · Decisiones de patrones documentadas", VIOLET),
    ("📚  Documentación", "16 documentos wiki · READMEs de 3 servicios\nSAD (Software Architecture Document) · Diagramas Mermaid", VIOLET),
    ("🎨  UX / UI", "Diseño de flujos de usuario · Mockups mobile\nComponentes shadcn/ui · Sistema de diseño y paleta", VIOLET),
    ("🔄  Refactorización", "Separación de módulos · Extracción de servicios\nAplicación de patrones clean architecture · Code review", VIOLET),
]
for i, (title, body, color) in enumerate(dev_ia):
    row, col = divmod(i, 3)
    x = 0.35 + col * 4.3
    y = 2.5 + row * 2.1
    card(s, x, y, 4.1, 1.9, title, body, color)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 · RUTA DEL DEMO
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "DEMO · RUTA DE DEMOSTRACIÓN", 22)
heading(s, "Demo en vivo: flujo completo en 8 pasos · ~7 minutos.")
subheading(s, "Si Gemini no responde, el sistema usa el motor de reglas — sin error visible.")

demo = [
    ("01", "📱 QR + APK\nInstalar y abrir"),
    ("02", "🔐 Login Paciente\nJWT + Auth0"),
    ("03", "🤖 Triage IA\nRed flags visibles"),
    ("04", "💳 Checkout\nPago simulado COP"),
    ("05", "👨‍⚕️ Panel Doctor\nCola prioridad ALTA"),
    ("06", "💬 Chat Real-time\nSocket.IO bidireccional"),
    ("07", "📝 Resumen IA\nGemini + Guardrail"),
    ("08", "📊 Dashboard\nKPIs reales BD"),
]
cols_demo = 4
for i, (num, step) in enumerate(demo):
    row, col = divmod(i, cols_demo)
    x = 0.4 + col * 3.25
    y = 2.5 + row * 1.9
    box(s, x, y, 3.1, 1.7, bg=PANEL2, border=BLUE)
    txt(s, num, x+0.1, y+0.1, 0.6, 0.5, size=20, bold=True, color=BLUE, font_name="Courier New")
    txt(s, step, x+0.75, y+0.15, 2.2, 1.4, size=11, bold=True, color=TEXT)

contingencias = [
    ("Gemini falla → RedFlagsEngine automático", CORAL),
    ("Problema de red → screenshots de respaldo", AMBER),
    ("Dashboard → datos reales de MongoDB", MINT),
]
for i, (msg, color) in enumerate(contingencias):
    x = 0.4 + i * 4.3
    box(s, x, 6.45, 4.1, 0.65, bg=PANEL, border=color)
    txt(s, msg, x+0.1, 6.52, 3.9, 0.55, size=10, bold=True, color=color)
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 · DEMO EN VIVO
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "DEMO EN VIVO", 23)
txt(s, "Demo activo — sistema funcionando en tiempo real.", 0.5, 0.8, 9.5, 1.4,
    size=36, bold=True, color=TEXT)
txt(s, "Backend desplegado · MongoDB Atlas con datos reales · Socket.IO activo",
    0.5, 2.35, 9.5, 0.5, size=14, color=MUTED)

live = [
    "Triage generado por Gemini 2.5-flash",
    "Chat en tiempo real via Socket.IO",
    "Resumen clínico con guardrail activo",
    "KPIs leídos directamente de MongoDB",
]
for i, item in enumerate(live):
    txt(s, f"→  {item}", 0.5, 3.0 + i*0.65, 9.0, 0.6, size=14, bold=True, color=TEXT)

# Big QR
box(s, 10.3, 1.2, 2.7, 3.8, bg=WHITE)
qr_path = ASSETS / "qr-apk.png"
if qr_path.exists():
    s.shapes.add_picture(str(qr_path), Inches(10.55), Inches(1.4), Inches(2.2), Inches(2.2))
txt(s, "APK Disponible", 10.3, 3.7, 2.7, 0.4, size=12, bold=True,
    color=RGBColor(0x0F,0x17,0x2A), align=PP_ALIGN.CENTER)
txt(s, "Expo Go / Descarga directa", 10.3, 4.1, 2.7, 0.35, size=9,
    color=RGBColor(0x47,0x55,0x69), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 · RESULTADOS
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "RESULTADOS LOGRADOS", 24)
heading(s, "Lo que construimos en números.")
subheading(s, "5 meses · 4 personas · 3 servicios productivos · 1 plataforma clínica funcionando.")

results = [
    ("Módulos NestJS backend", "17", BLUE),
    ("Tests pasando (0 fallando)", "406+", MINT),
    ("Cobertura statements", "93%", VIOLET),
    ("Umbral mínimo CI (branches)", "80%", AMBER),
    ("Suites E2E con MongoDB real", "12", BLUE),
    ("Roles de usuario diferenciados", "3", MINT),
    ("Frontends independientes", "2", VIOLET),
    ("Especialidades médicas", "3", AMBER),
]
for i, (key, val, color) in enumerate(results):
    row, col = divmod(i, 2)
    x = 0.4 + col * 6.5
    y = 2.5 + row * 1.1
    box(s, x, y, 6.2, 0.95, bg=PANEL, border=color)
    txt(s, key, x+0.2, y+0.1, 4.5, 0.75, size=12, color=TEXT)
    txt(s, val, x+4.8, y+0.1, 1.2, 0.75, size=16, bold=True, color=color,
        align=PP_ALIGN.RIGHT, font_name="Courier New")
qr_badge(s)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 · CONCLUSIONES
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
tag(s, "CONCLUSIONES", 25)
heading(s, "Un sistema que funciona, con tests que lo respaldan y arquitectura que se puede defender.")
subheading(s, "SaludDeUna resuelve un problema real de comunicación clínica con tecnología de producción.")

learnings = [
    "La IA en salud no puede ser un black box — necesita guardrails explícitos en código",
    "Tests E2E con BD real son críticos — los mocks no detectan bugs de migración",
    "La observabilidad desde el día 1 ahorra horas de debugging en producción",
    "La modularidad del backend permitió desarrollar en paralelo sin conflictos",
]
txt(s, "Lo que aprendimos:", 0.5, 2.3, 8.5, 0.4, size=13, bold=True, color=MUTED)
for i, learn in enumerate(learnings):
    txt(s, f"→  {learn}", 0.5, 2.8 + i*0.7, 8.5, 0.65, size=13, color=TEXT)

# QR + preguntas
box(s, 9.5, 2.0, 3.5, 4.8, bg=WHITE)
if qr_path.exists():
    s.shapes.add_picture(str(qr_path), Inches(9.75), Inches(2.2), Inches(3.0), Inches(3.0))
txt(s, "APK disponible", 9.5, 5.3, 3.5, 0.4, size=12, bold=True,
    color=RGBColor(0x0F,0x17,0x2A), align=PP_ALIGN.CENTER)

txt(s, "¿Preguntas?", 0.5, 5.8, 8.5, 1.0, size=48, bold=True, color=BLUE)
txt(s, "El código es abierto. Los tests están pasando. La demo acaba de funcionar.",
    0.5, 6.75, 8.5, 0.45, size=11, color=MUTED)

# ═══════════════════════════════════════════════════════════════════════════
# BACKUP SLIDES
# ═══════════════════════════════════════════════════════════════════════════

def backup_slide(title_text, subtitle_text, image_file, num_label):
    s = new_slide()
    box(s, 0, 0, 13.33, 7.5, bg=RGBColor(0x07,0x0E,0x26))
    tag(s, f"BACKUP · {num_label}", 0)
    heading(s, title_text)
    subheading(s, subtitle_text)
    img(s, image_file, 0.5, 2.3, 12.33, 4.6)
    qr_badge(s)

backup_slide("C4 Nivel 1 — Contexto del Sistema", "3 actores · 1 sistema central · 6 servicios externos",
             "sad_c4_context_diagram.png", "C4 CONTEXTO")
backup_slide("C4 Nivel 2 — Vista de Contenedores", "API + Worker + Mobile + Web + MongoDB + Redis",
             "sad_c4_container_view.png", "C4 CONTENEDORES")
backup_slide("Despliegue Objetivo AWS — ECS Fargate", "ECS Fargate SPOT · ALB · Route53 · CloudFront · WAF (~$25-31/mes)",
             "sad_aws_deployment_diagram.png", "AWS DEPLOYMENT")

# Data models (two diagrams)
s = new_slide()
box(s, 0, 0, 13.33, 7.5, bg=RGBColor(0x07,0x0E,0x26))
tag(s, "BACKUP · MODELO DE DATOS", 0)
heading(s, "Modelo de datos: dominio usuarios y dominio clínico.")
subheading(s, "MongoDB + Mongoose · Atlas Vector Search HNSW 768d para RAG")
img(s, "sad_data_model_users.png", 0.3, 2.3, 6.2, 4.5)
img(s, "sad_data_model_clinical.png", 6.9, 2.3, 6.2, 4.5)
qr_badge(s)

# CI/CD + Observabilidad
s = new_slide()
box(s, 0, 0, 13.33, 7.5, bg=RGBColor(0x07,0x0E,0x26))
tag(s, "BACKUP · CI/CD Y OBSERVABILIDAD", 0)
heading(s, "CI/CD y observabilidad reducen el riesgo operacional.")
subheading(s, "GitHub Actions · lint → build → test → E2E · OpenTelemetry · logs JSON")
img(s, "sad_cicd_pipeline.png", 0.3, 2.3, 6.2, 4.5)
img(s, "sad_observability_stack.png", 6.9, 2.3, 6.2, 4.5)
qr_badge(s)

# Q&A
s = new_slide()
box(s, 0, 0, 13.33, 7.5, bg=RGBColor(0x07,0x0E,0x26))
tag(s, "BACKUP · PREGUNTAS DEL JURADO", 0)
heading(s, "Preguntas técnicas probables y respuestas preparadas.")

qa = [
    ("¿Cómo garantizan que la IA no diagnostica?",
     "El guardrail es código, no promesa. Si Gemini genera lenguaje de diagnóstico, aiSummary = null. Podemos mostrar GuardrailService.", VIOLET),
    ("¿Por qué MongoDB y no PostgreSQL?",
     "Vector Search nativo para RAG, esquema flexible para datos clínicos que evolucionan, embeddings 768d nativos en Atlas.", MINT),
    ("¿Por qué NestJS y no Spring Boot?",
     "DI explícita, módulos aislados, TypeScript end-to-end, ecosistema para WebSocket y RAG. Sin overhead de Spring.", BLUE),
    ("¿Cómo escala el chat?",
     "Redis adapter para Socket.IO. En multi-instancia, pub/sub distribuye mensajes entre instancias sin estado compartido.", AMBER),
    ("¿Qué pasa si Gemini falla?",
     "RedFlagsEngine basado en reglas locales toma el control. analysisType = RULE_BASED. Sin error visible para el usuario.", CORAL),
    ("¿Costo de infraestructura?",
     "Dev: Railway + Vercel + Atlas M0 + Redis free ≈ $0/mes. Objetivo producción: ECS Fargate SPOT ≈ $25-31/mes.", BLUE),
]
for i, (q, a, color) in enumerate(qa):
    row, col = divmod(i, 2)
    x = 0.35 + col * 6.5
    y = 2.5 + row * 1.6
    box(s, x, y, 6.2, 1.45, bg=PANEL2, border=color)
    txt(s, f"Q: {q}", x+0.15, y+0.1, 5.9, 0.5, size=10, bold=True, color=color)
    txt(s, a, x+0.15, y+0.62, 5.9, 0.75, size=10, color=TEXT)
qr_badge(s)

# Business Model Canvas
s = new_slide()
box(s, 0, 0, 13.33, 7.5, bg=RGBColor(0x07,0x0E,0x26))
tag(s, "BACKUP · BUSINESS MODEL CANVAS", 0)
heading(s, "Business Model Canvas y viabilidad de mercado.")
img(s, "Business-Model-Canvas.png", 0.4, 2.2, 12.5, 4.8)
qr_badge(s)

# ── Save ─────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"✅  PPTX generado: {OUT}")
print(f"   Diapositivas: {len(prs.slides)}")
